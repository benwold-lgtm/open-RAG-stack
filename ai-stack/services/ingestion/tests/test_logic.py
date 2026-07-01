"""Unit tests for ingestion's core logic — the deterministic transforms behind the
ingest endpoints: text chunking, whitespace cleaning, multi-format document
extraction, and file storage. No Qdrant / embedding service involved.

Async helpers are driven with asyncio.run() so no asyncio pytest plugin is needed.
(Importing `main` pulls the full ingestion stack; run where its requirements are installed.)"""
import asyncio
import io

import pytest

import main


# ── chunk_text ────────────────────────────────────────────────────────────────
def test_chunk_text_short_text_is_one_chunk():
    assert main.chunk_text("hello world", chunk_size=100, overlap=10) == ["hello world"]


def test_chunk_text_splits_long_text_into_many():
    text = " ".join(f"word{i}" for i in range(500))
    chunks = main.chunk_text(text, chunk_size=100, overlap=20)
    assert len(chunks) > 1
    assert all(c.strip() for c in chunks)            # no empty/whitespace-only chunks
    assert any("word499" in c for c in chunks)       # the tail of the input is covered


def test_chunk_text_whitespace_only_yields_nothing():
    assert main.chunk_text("   \n\n  ", chunk_size=50, overlap=5) == []


# ── _clean ────────────────────────────────────────────────────────────────────
def test_clean_collapses_runs_of_spaces_and_blank_lines():
    assert main._clean("a    b\n\n\n\nc  ") == "a b\n\nc"


# ── extract_document ──────────────────────────────────────────────────────────
def test_extract_document_txt_cleans_text():
    title, segs = asyncio.run(main.extract_document("notes.txt", b"line one\n\n\n\nline two   "))
    assert title == "notes.txt"
    assert len(segs) == 1 and segs[0]["page"] is None and segs[0]["has_image"] is False
    assert segs[0]["text"] == "line one\n\nline two"


def test_extract_document_md_decodes_utf8():
    _, segs = asyncio.run(main.extract_document("r.md", "café résumé".encode("utf-8")))
    assert "café" in segs[0]["text"] and "résumé" in segs[0]["text"]


def test_extract_document_rejects_unsupported_extension():
    with pytest.raises(ValueError):
        asyncio.run(main.extract_document("malware.exe", b"x"))


def test_extract_document_docx():
    from docx import Document
    d = Document()
    d.add_paragraph("Hello world")
    d.add_paragraph("Second paragraph")
    buf = io.BytesIO()
    d.save(buf)
    _, segs = asyncio.run(main.extract_document("x.docx", buf.getvalue()))
    assert "Hello world" in segs[0]["text"] and "Second paragraph" in segs[0]["text"]


def test_extract_document_pptx():
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])   # "Title Only" layout
    slide.shapes.title.text = "Slide Heading"
    buf = io.BytesIO()
    prs.save(buf)
    _, segs = asyncio.run(main.extract_document("x.pptx", buf.getvalue()))
    assert "Slide Heading" in segs[0]["text"]


# ── save_file ─────────────────────────────────────────────────────────────────
def test_save_file_legacy_writes_to_files_dir(tmp_path, monkeypatch):
    # No DOC_STORE -> legacy FILES_DIR/<doc_id>.<ext>, returns None (no relative path).
    monkeypatch.setattr(main, "DOC_STORE", "")
    monkeypatch.setattr(main, "FILES_DIR", str(tmp_path))
    assert main.save_file("doc123", "report.PDF", b"%PDF-bytes") is None
    assert (tmp_path / "doc123.pdf").read_bytes() == b"%PDF-bytes"


def test_save_file_doc_store_by_vendor(tmp_path, monkeypatch):
    monkeypatch.setattr(main, "DOC_STORE", str(tmp_path))
    rel = main.save_file("docabc", "PowerEdge Guide.pdf", b"data", vendor="Dell")
    assert rel == "Dell/PowerEdge Guide.pdf"
    assert (tmp_path / "Dell" / "PowerEdge Guide.pdf").read_bytes() == b"data"


def test_save_file_collision_suffixes_doc_id(tmp_path, monkeypatch):
    monkeypatch.setattr(main, "DOC_STORE", str(tmp_path))
    main.save_file("aaaaaaaaaaaaaaaa", "x.pdf", b"first", vendor="hpe")
    rel = main.save_file("bbbbbbbbbbbbbbbb", "x.pdf", b"second", vendor="hpe")  # same name, diff bytes
    assert rel == "hpe/x-bbbbbbbb.pdf"
    assert (tmp_path / "hpe" / "x.pdf").read_bytes() == b"first"
    assert (tmp_path / "hpe" / "x-bbbbbbbb.pdf").read_bytes() == b"second"


def test_save_file_path_traversal_is_contained(tmp_path, monkeypatch):
    monkeypatch.setattr(main, "DOC_STORE", str(tmp_path))
    rel = main.save_file("docx0000", "../../evil.pdf", b"x", vendor="../etc")
    assert rel == "etc/evil.pdf"                      # basename-only sanitisation


# ── figure-caption extraction (isolated, document-native) ─────────────────────
def test_captions_joins_split_label_and_sentence():
    # PyMuPDF often emits the label ("Figure 1.") and its caption sentence as separate blocks.
    texts = ["Figure 1.",
             "High-level diagram of Kubernetes and RAG components with Dell infrastructure",
             "Dell AI Factory", "The Dell AI Factory platform is partnered with vendors."]
    assert main._captions_from_block_texts(texts) == [
        "Figure 1. High-level diagram of Kubernetes and RAG components with Dell infrastructure"]


def test_captions_inline_single_block():
    assert main._captions_from_block_texts(["Figure 2. Reference architecture overview"]) == [
        "Figure 2. Reference architecture overview"]


def test_captions_table_and_dedup():
    texts = ["Table 1. Comparison of options", "Table 1. Comparison of options", "body text"]
    assert main._captions_from_block_texts(texts) == ["Table 1. Comparison of options"]


def test_captions_drops_bare_label_without_wording():
    # "Figure 3." has no caption sentence (the next block is itself a caption) -> dropped.
    assert main._captions_from_block_texts(["Figure 3.", "Figure 4. Real caption text"]) == [
        "Figure 4. Real caption text"]


def test_captions_short_complete_caption_not_joined_with_prose():
    # A short but complete caption must NOT absorb the following prose block.
    texts = ["Figure 5. Topology", "This paragraph is unrelated prose."]
    assert main._captions_from_block_texts(texts) == ["Figure 5. Topology"]


def test_captions_none_in_prose():
    assert main._captions_from_block_texts(["Just some prose about RAG and Dell servers."]) == []


# ── caption title-enrichment (findable by the document's topic) ───────────────
def test_enrich_caption_prefixes_document_title():
    # The document's own title carries the topic ("Redis") the bare caption lacks.
    assert main.enrich_caption(
        "Bring Real-Time RAG into Production with Dell AI Factory and Redis",
        "Figure 1. High-level diagram of Kubernetes and RAG components",
    ) == ("Bring Real-Time RAG into Production with Dell AI Factory and Redis. "
          "Figure 1. High-level diagram of Kubernetes and RAG components")


def test_enrich_caption_noop_without_title():
    assert main.enrich_caption("", "Figure 2. Reference architecture") == "Figure 2. Reference architecture"


class _Doc:
    def __init__(self, title):
        self.metadata = {"title": title}


def test_pdf_doc_title_keeps_real_title():
    assert main._pdf_doc_title(_Doc("Dell AI Factory with Redis Reference Design")) == \
        "Dell AI Factory with Redis Reference Design"


def test_pdf_doc_title_drops_generic_and_filename_like():
    assert main._pdf_doc_title(_Doc("untitled")) == ""
    assert main._pdf_doc_title(_Doc("PowerPoint Presentation")) == ""
    assert main._pdf_doc_title(_Doc("h20163-genAI-Redis-RAG.drd.pdf")) == ""
    assert main._pdf_doc_title(_Doc("")) == ""
    assert main._pdf_doc_title(_Doc(None)) == ""
