import os
import re
import shutil
import hashlib
import asyncio
import aiosqlite
import httpx
from langchain_text_splitters import RecursiveCharacterTextSplitter
from datetime import datetime
from fastapi import FastAPI, HTTPException, BackgroundTasks, UploadFile, File, Form, Depends
from pydantic import BaseModel
from typing import Optional
from crawl4ai import AsyncWebCrawler, CrawlerRunConfig, CacheMode
from qdrant_client import QdrantClient
from qdrant_client.models import (
    Distance, VectorParams, PointStruct, Filter,
    FieldCondition, MatchValue
)
from tenacity import retry, stop_after_attempt, wait_exponential
from rag_auth import RoleScopes, Principal, AUTH_LOCAL
from rag_auth.authenticator import StaticKeyAuthenticator
from rag_auth.deps import make_authenticate_request, require_scope, verify_boot_config
from prometheus_fastapi_instrumentator import Instrumentator

app = FastAPI(title="Ingestion Service")

# Prometheus metrics at GET /metrics — open like /health (only request counts/latencies, no
# secrets). Network-gate it the same way as the rest of the data plane.
Instrumentator().instrument(app).expose(app)

# ── Authentication (shared service token via rag_auth) ────────────────────────
# Callers (rag-admin, ai-agent's lexical search, the helper scripts) present a
# shared SERVICE_TOKEN as a bearer credential. Reads need ingest:read, writes
# ingest:write, destructive ops docs:manage. Fail-closed in production unless an
# operator explicitly opts into anonymous access on a trusted, isolated LAN.
# NOTE: the page-image and document-file routes stay open — browsers load those URLs
# directly (see INGESTION_PUBLIC_URL) and cannot send a bearer token. Gate at the network.
ENVIRONMENT     = os.getenv("ENVIRONMENT",     "development")
SERVICE_TOKEN   = os.getenv("SERVICE_TOKEN",   "")
ALLOW_ANONYMOUS = os.getenv("ALLOW_ANONYMOUS", "false").lower() in ("1", "true", "yes")

INGEST_READ, INGEST_WRITE, DOCS_MANAGE = "ingest:read", "ingest:write", "docs:manage"
_ROLES = RoleScopes({
    "admin":  {INGEST_READ, INGEST_WRITE, DOCS_MANAGE},
    "viewer": {INGEST_READ},
})
_keys: dict[str, Principal] = {}
if SERVICE_TOKEN:
    _keys[SERVICE_TOKEN] = Principal(
        subject="key:service", scopes=_ROLES.scopes_for_role("admin"), auth_method=AUTH_LOCAL
    )
_authn = make_authenticate_request(StaticKeyAuthenticator(_keys))

for _w in verify_boot_config(
    production=ENVIRONMENT == "production",
    any_auth_enabled=bool(SERVICE_TOKEN),
    allow_anonymous=ALLOW_ANONYMOUS,
):
    print(f"[ingestion] WARN: {_w}", flush=True)

# Enforce auth only when a token is configured AND anonymous access isn't allowed.
_AUTH_ACTIVE = bool(SERVICE_TOKEN) and not ALLOW_ANONYMOUS

def _guard(scope: str):
    return [Depends(_authn), Depends(require_scope(scope))] if _AUTH_ACTIVE else []

crawl_semaphore = asyncio.Semaphore(1)
embed_semaphore = asyncio.Semaphore(2)

# ── Configuration ─────────────────────────────────────────────────────────────
QDRANT_URL          = os.getenv("QDRANT_URL",          "http://qdrant.qdrant.svc.cluster.local:6333")
QDRANT_API_KEY      = os.getenv("QDRANT_API_KEY")
EMBEDDING_URL       = os.getenv("EMBEDDING_URL",       "http://embedding.embedding.svc.cluster.local:8001")
DB_PATH             = os.getenv("DB_PATH",             "/app/data/ingestion.db")
FILES_DIR           = os.getenv("FILES_DIR",           "/app/data/files")
# Optional shared document store (e.g. an NFS mount). When set, original uploads are saved as
# <DOC_STORE>/<vendor>/<filename> (browsable by vendor) and served back via /documents/{id}/file
# so chat citations can link to the source. Empty = keep the legacy FILES_DIR/<doc_id>.<ext>.
DOC_STORE           = os.getenv("DOC_STORE",           "")
WATCH_DIR           = os.getenv("WATCH_DIR",           "")
WATCH_POLL_INTERVAL = int(os.getenv("WATCH_POLL_INTERVAL", "60"))
CHUNK_SIZE               = int(os.getenv("CHUNK_SIZE",          "512"))
CHUNK_OVERLAP            = int(os.getenv("CHUNK_OVERLAP",       "50"))
EMBEDDING_DIM            = int(os.getenv("EMBEDDING_DIM",       "768"))

# OCR fallback for image-rich PDF pages (Tesseract, CPU-only). A page is rendered
# and OCR'd only when its text layer is sparse — born-digital text pages never hit it.
OCR_ENABLED   = os.getenv("OCR_ENABLED", "true").lower() in ("1", "true", "yes")
OCR_MIN_CHARS = int(os.getenv("OCR_MIN_CHARS", "100"))
OCR_DPI       = int(os.getenv("OCR_DPI",       "200"))
# A page with more than this many vector-drawing ops is treated as having a diagram
# worth surfacing as an image (distinguishes a diagram from a stray underline/rule).
DRAWING_IMG_THRESHOLD = 5

SUPPORTED_EXTENSIONS = {"pdf", "docx", "pptx", "txt", "md"}

# ── Database setup ────────────────────────────────────────────────────────────
async def init_db():
    async with aiosqlite.connect(DB_PATH) as db:
        await db.execute("""
            CREATE VIRTUAL TABLE IF NOT EXISTS fts USING fts5(
                point_id UNINDEXED, collection UNINDEXED, url UNINDEXED,
                content, title, vendor, doc_id UNINDEXED
            )
        """)
        await db.execute("""
            CREATE TABLE IF NOT EXISTS documents (
                id           TEXT PRIMARY KEY,
                url          TEXT NOT NULL,
                collection   TEXT NOT NULL,
                vendor       TEXT,
                title        TEXT,
                content_hash TEXT,
                status       TEXT DEFAULT 'pending',
                chunk_count  INTEGER DEFAULT 0,
                error        TEXT,
                source_type  TEXT DEFAULT 'url',
                created_at   TEXT,
                updated_at   TEXT,
                last_checked TEXT
            )
        """)
        # Migration: add source_type to tables created before this column existed
        try:
            await db.execute("ALTER TABLE documents ADD COLUMN source_type TEXT DEFAULT 'url'")
        except Exception:
            pass
        # Migration: file_path = path of the stored original within DOC_STORE (vendor/filename),
        # used to serve /documents/{id}/file. NULL for url docs and legacy FILES_DIR-only docs.
        try:
            await db.execute("ALTER TABLE documents ADD COLUMN file_path TEXT")
        except Exception:
            pass
        await db.commit()

async def _probe_embedding_dim():
    """Warn loudly at startup if EMBEDDING_DIM doesn't match the live embedding service."""
    try:
        async with httpx.AsyncClient(timeout=10.0) as client:
            resp = await client.post(f"{EMBEDDING_URL}/v1/embeddings", json={"input": ["dim-probe"]})
            resp.raise_for_status()
            actual = len(resp.json()["data"][0]["embedding"])
        if actual != EMBEDDING_DIM:
            print(
                f"[ERROR] Embedding dimension mismatch: service returned {actual}, "
                f"EMBEDDING_DIM={EMBEDDING_DIM}. "
                f"Update config.embeddingDim in ingestion/values.yaml to {actual} "
                f"and re-ingest all documents — searches will silently return nothing until fixed.",
                flush=True,
            )
    except Exception as exc:
        print(f"[WARN] Could not probe embedding service at startup ({exc}) — skipping dim check.", flush=True)


@app.on_event("startup")
async def startup():
    os.makedirs(os.path.dirname(DB_PATH), exist_ok=True)
    os.makedirs(FILES_DIR, exist_ok=True)
    await init_db()
    await _probe_embedding_dim()
    if WATCH_DIR:
        asyncio.create_task(watch_folder())

# ── Qdrant client ─────────────────────────────────────────────────────────────
def get_qdrant():
    return QdrantClient(url=QDRANT_URL, api_key=QDRANT_API_KEY)

# ── Ensure collection exists ──────────────────────────────────────────────────
async def ensure_collection(collection: str):
    client = get_qdrant()
    existing = [c.name for c in client.get_collections().collections]
    if collection not in existing:
        client.create_collection(
            collection_name=collection,
            vectors_config=VectorParams(size=EMBEDDING_DIM, distance=Distance.COSINE)
        )

def _copy_points(client, source: str, target: str, scroll_filter=None, batch: int = 256) -> int:
    """Copy points (vectors + payload) from one Qdrant collection to another,
    re-tagging each payload's `collection` field. Point IDs are preserved so the
    FTS5 rows (keyed by point_id) stay valid. Returns the number of points copied.
    All collections share one vector config (EMBEDDING_DIM), so no re-embedding."""
    copied = 0
    offset = None
    while True:
        points, offset = client.scroll(
            collection_name=source,
            scroll_filter=scroll_filter,
            with_payload=True,
            with_vectors=True,
            limit=batch,
            offset=offset,
        )
        if not points:
            break
        retagged = [
            PointStruct(id=p.id, vector=p.vector, payload={**(p.payload or {}), "collection": target})
            for p in points
        ]
        client.upsert(collection_name=target, points=retagged)
        copied += len(retagged)
        if offset is None:
            break
    return copied

# ── Text chunking ─────────────────────────────────────────────────────────────
def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=overlap)
    return [c for c in splitter.split_text(text) if c.strip()]

# ── Web fetcher (Crawl4AI) ────────────────────────────────────────────────────
@retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=2, max=10))
async def fetch_url(url: str) -> tuple[str, str]:
    """Fetch a URL using a headless browser and return (title, clean_text)."""
    async with crawl_semaphore:
        base_config = dict(
            cache_mode=CacheMode.BYPASS,
            wait_until="domcontentloaded",
            page_timeout=45000,
            remove_overlay_elements=True,
            excluded_tags=["nav", "footer", "header", "aside"],
            word_count_threshold=10,
            magic=True,
        )
        async with AsyncWebCrawler(headless=True) as crawler:
            result = await crawler.arun(
                url=url,
                config=CrawlerRunConfig(
                    **base_config,
                    css_selector="article, main, .blog-content, .article-body, #mh-main",
                )
            )

            # Retry without CSS selector if the selector matched nothing useful
            if result.success and (not result.markdown or len(result.markdown.strip()) < 100):
                result = await crawler.arun(
                    url=url, config=CrawlerRunConfig(**base_config)
                )

    if not result.success:
        raise ValueError(f"Failed to crawl {url}: {result.error_message}")

    title = (result.metadata or {}).get("title", "").strip()
    text = re.sub(r'\s+', ' ', result.markdown or "").strip()
    return title, text

# ── Document extractor ────────────────────────────────────────────────────────
def _ocr_page(page) -> str:
    """Render a PDF page to an image and OCR it with Tesseract (CPU-only)."""
    import io
    import pytesseract
    from PIL import Image
    pix = page.get_pixmap(dpi=OCR_DPI)
    img = Image.open(io.BytesIO(pix.tobytes("png")))
    return pytesseract.image_to_string(img)


_CAPTION_RE = re.compile(r'^(figure|fig\.?|table|diagram|exhibit)\s+\d+', re.IGNORECASE)


def _captions_from_block_texts(texts: list[str]) -> list[str]:
    """Pick figure/table captions out of a page's text blocks (in reading order).

    Pure/testable core of _extract_captions: a caption is the document's OWN words describing a
    figure ("Figure 3. <text>"). A bare label block ("Figure 1.") is joined with the caption
    sentence that follows it; bare labels with no real wording are dropped."""
    captions: list[str] = []
    for idx, t in enumerate(texts):
        if not _CAPTION_RE.match(t):
            continue
        cap = t
        # A bare label block ("Figure 1." with nothing after the number) sits just above its
        # caption sentence — join them. A short-but-complete caption is left untouched.
        label_only = len(_CAPTION_RE.sub("", t, count=1).strip(" .:-")) < 3
        if label_only and idx + 1 < len(texts) and not _CAPTION_RE.match(texts[idx + 1]):
            cap = f"{t} {texts[idx + 1]}".strip()
        cap = cap[:300]
        remainder = _CAPTION_RE.sub("", cap, count=1).strip(" .:-")
        if len(remainder) >= 5 and cap not in captions:    # keep only captions with real wording
            captions.append(cap)
    return captions


def _extract_captions(page) -> list[str]:
    """Pull figure/table captions from a PDF page as clean, standalone strings, so each can be
    indexed as its own retrieval unit (instead of buried in a prose chunk) — making a diagram
    findable by the surrounding context a user would actually search for."""
    blocks = [b for b in page.get_text("blocks") if len(b) >= 5 and (len(b) < 7 or b[6] == 0)]
    blocks.sort(key=lambda b: (round(b[1]), round(b[0])))   # reading order: top-down, left-right
    texts = [" ".join((b[4] or "").split()) for b in blocks]
    return _captions_from_block_texts(texts)


_GENERIC_TITLES = {"untitled", "title", "powerpoint presentation", "document", "presentation"}


def _pdf_doc_title(doc) -> str:
    """The document's own embedded title (PDF metadata), used to enrich figure-caption points.

    A caption often only describes the figure ("Figure 1. High-level diagram of Kubernetes and
    RAG components") without naming the product the document is about, so it isn't found by a
    topic query like "Dell RAG Redis diagram". Prefixing the caption with the document's own
    title — its real, author-written words, not an AI interpretation — restores that topic
    context. Returns '' when the metadata title is missing or too generic to help."""
    t = ((doc.metadata or {}).get("title") or "").strip()
    if len(t) < 5 or t.lower() in _GENERIC_TITLES or t.lower().endswith(".pdf"):
        return ""
    return t[:160]


def enrich_caption(doc_title: str, caption: str) -> str:
    """Prefix a caption with the document's title so the figure is findable by the document's
    topic. No-op when the title is empty (filename-style titles embed poorly — skip them)."""
    return f"{doc_title}. {caption}" if doc_title else caption


def _extract_pdf_pages(content: bytes) -> list[dict]:
    """Return per-page records [{page, text, has_image, captions}] using PyMuPDF.

    Pages with a sparse text layer are rendered and OCR'd so the labels in
    image/diagram pages (device names, IPs, VLANs) become searchable. has_image
    flags pages carrying a raster image or a diagram for the UI to surface."""
    import fitz  # PyMuPDF
    pages = []
    with fitz.open(stream=content, filetype="pdf") as doc:
        doc_title = _pdf_doc_title(doc)
        for i, page in enumerate(doc, start=1):
            text = page.get_text() or ""
            has_raster = bool(page.get_images())
            ocr_applied = False
            if OCR_ENABLED and len(text.strip()) < OCR_MIN_CHARS:
                ocr_text = _ocr_page(page)
                if ocr_text.strip():
                    text = ocr_text
                    ocr_applied = True
            has_diagram = len(page.get_drawings()) > DRAWING_IMG_THRESHOLD
            pages.append({
                "page": i,
                "text": text,
                "has_image": has_raster or has_diagram or ocr_applied,
                "captions": [enrich_caption(doc_title, c) for c in _extract_captions(page)],
            })
    return pages


def _extract_docx(content: bytes) -> str:
    import io
    from docx import Document
    doc = Document(io.BytesIO(content))
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(content: bytes) -> str:
    import io
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    slides = []
    for slide in prs.slides:
        slide_text = "\n".join(
            shape.text for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        )
        if slide_text.strip():
            slides.append(slide_text)
    return "\n\n".join(slides)


def _office_to_pdf(content: bytes, ext: str) -> bytes:
    """Render a DOCX/PPTX to PDF with headless LibreOffice, so office docs flow through the same
    page-aware pipeline as PDFs (per-page segments, figure captions, OCR, page-image rendering).
    A private per-call UserInstallation profile keeps concurrent conversions from colliding.
    Raises on failure so the caller can fall back to plain text extraction."""
    import subprocess, tempfile, glob
    with tempfile.TemporaryDirectory() as tmp:
        src = os.path.join(tmp, f"input.{ext}")
        with open(src, "wb") as f:
            f.write(content)
        subprocess.run(
            ["soffice", "--headless", f"-env:UserInstallation=file://{tmp}/lo",
             "--convert-to", "pdf", "--outdir", tmp, src],
            check=True, capture_output=True, timeout=180,
            env={**os.environ, "HOME": tmp},
        )
        pdfs = glob.glob(os.path.join(tmp, "*.pdf"))
        if not pdfs:
            raise RuntimeError("LibreOffice produced no PDF output")
        with open(pdfs[0], "rb") as f:
            return f.read()


def _clean(text: str) -> str:
    text = re.sub(r' {2,}', ' ', text)
    return re.sub(r'\n{3,}', '\n\n', text).strip()


async def extract_document(filename: str, content: bytes) -> tuple[str, list[dict], Optional[bytes]]:
    """Extract (title, segments, viewable_pdf) from a PDF, DOCX, PPTX, txt, or md file.

    Each segment is {page, text, has_image}. PDFs and (via LibreOffice) DOCX/PPTX yield one
    segment per page — page-aware, OCR fallback, figure captions. viewable_pdf is the page-accurate
    PDF to store/serve for page citations: the file itself for a PDF, the LibreOffice-derived PDF
    for an office doc, or None for txt/md and when office->PDF conversion is unavailable/fails."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: .{ext}. Supported: {', '.join(SUPPORTED_EXTENSIONS)}")

    loop = asyncio.get_event_loop()
    viewable: Optional[bytes] = None
    if ext == "pdf":
        segments = await loop.run_in_executor(None, _extract_pdf_pages, content)
        viewable = content
    elif ext in ("docx", "pptx"):
        try:
            pdf_bytes = await loop.run_in_executor(None, _office_to_pdf, content, ext)
            segments = await loop.run_in_executor(None, _extract_pdf_pages, pdf_bytes)
            viewable = pdf_bytes
        except Exception as exc:                                  # LibreOffice missing/failed
            print(f"[ingestion] {ext}->PDF conversion failed for {filename} ({exc}); "
                  f"falling back to text-only extraction.", flush=True)
            extractor = _extract_docx if ext == "docx" else _extract_pptx
            text = await loop.run_in_executor(None, extractor, content)
            segments = [{"page": None, "text": text, "has_image": False}]
    else:
        text = content.decode("utf-8", errors="replace")
        segments = [{"page": None, "text": text, "has_image": False}]

    for seg in segments:
        seg["text"] = _clean(seg["text"])
    return filename, segments, viewable

# ── File storage ──────────────────────────────────────────────────────────────
def _safe_component(name: str) -> str:
    """basename only (no path traversal), leading dots stripped."""
    return os.path.basename((name or "").strip()).lstrip(".")


def save_file(doc_id: str, filename: str, content: bytes, vendor: str = "") -> Optional[str]:
    """Persist the original upload. With DOC_STORE set, write to <store>/<vendor>/<filename>
    (browsable by vendor) and return that RELATIVE path ('vendor/filename'); on a name clash with
    a different file, the name is suffixed with a short doc_id. Without DOC_STORE, fall back to
    the legacy FILES_DIR/<doc_id>.<ext> and return None."""
    if DOC_STORE:
        safe_vendor = _safe_component(vendor) or "unknown"
        safe_name = _safe_component(filename) or f"{doc_id}.bin"
        os.makedirs(os.path.join(DOC_STORE, safe_vendor), exist_ok=True)
        target = os.path.join(DOC_STORE, safe_vendor, safe_name)
        if os.path.exists(target):
            with open(target, "rb") as existing:
                if existing.read() != content:                    # same name, different file
                    stem, ext = os.path.splitext(safe_name)
                    safe_name = f"{stem}-{doc_id[:8]}{ext}"
                    target = os.path.join(DOC_STORE, safe_vendor, safe_name)
        with open(target, "wb") as f:
            f.write(content)
        return f"{safe_vendor}/{safe_name}"

    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "bin"
    with open(os.path.join(FILES_DIR, f"{doc_id}.{ext}"), "wb") as f:
        f.write(content)
    return None


async def _locate_source(doc_id: str) -> Optional[str]:
    """Absolute path to a document's stored original — the DOC_STORE copy if recorded,
    otherwise the legacy FILES_DIR/<doc_id>.<ext>. None if neither exists."""
    if DOC_STORE:
        async with aiosqlite.connect(DB_PATH) as db:
            async with db.execute("SELECT file_path FROM documents WHERE id=?", (doc_id,)) as cur:
                row = await cur.fetchone()
        if row and row[0]:
            p = os.path.join(DOC_STORE, row[0])
            if os.path.exists(p):
                return p
    import glob
    matches = glob.glob(os.path.join(FILES_DIR, f"{doc_id}.*"))
    return matches[0] if matches else None

# ── Embedding caller ──────────────────────────────────────────────────────────
async def embed_texts(texts: list[str]) -> list[list[float]]:
    async with httpx.AsyncClient(timeout=120.0) as client:
        response = await client.post(
            f"{EMBEDDING_URL}/v1/embeddings",
            json={"input": texts}
        )
        response.raise_for_status()
        data = response.json()
        return [item["embedding"] for item in data["data"]]

# ── Shared pipeline: chunk → embed → upsert ───────────────────────────────────
async def run_pipeline(
    doc_id: str,
    source: str,
    title: str,
    segments: list[dict],
    content_hash: str,
    collection: str,
    vendor: str,
    access_roles: list[str],
    classification: str,
    source_type: str,
) -> int:
    """Chunk, embed, and upsert to Qdrant. Returns chunk count.

    segments is a list of {page, text, has_image}; each is chunked independently
    so every chunk carries its source page and image flag."""
    now = datetime.utcnow().isoformat()

    # (content, page, has_image, kind) — normal text chunks, plus each figure caption as its own
    # clean retrieval unit ("caption") so a diagram is findable by its surrounding wording.
    records = [
        (c, seg.get("page"), seg.get("has_image", False), "text")
        for seg in segments
        for c in chunk_text(seg["text"])
    ]
    records += [
        (cap, seg.get("page"), True, "caption")
        for seg in segments
        for cap in seg.get("captions", [])
    ]
    if not records:
        raise ValueError("No chunks generated from content")
    text_count = sum(1 for r in records if r[3] == "text")

    await ensure_collection(collection)

    async with aiosqlite.connect(DB_PATH) as db:
        await db.execute("DELETE FROM fts WHERE doc_id=?", (doc_id,))
        await db.commit()

    client = get_qdrant()
    client.delete(
        collection_name=collection,
        points_selector=Filter(
            must=[FieldCondition(key="doc_id", match=MatchValue(value=doc_id))]
        )
    )

    total_chunks = len(records)
    BATCH_SIZE = 8
    for i in range(0, total_chunks, BATCH_SIZE):
        batch = records[i:i + BATCH_SIZE]
        batch_embeddings = await embed_texts([rec[0] for rec in batch])

        points = []
        fts_rows = []
        for j, ((content, page, has_image, kind), embedding) in enumerate(zip(batch, batch_embeddings)):
            point_id = int.from_bytes(
                hashlib.sha256(f"{doc_id}-{i+j}".encode()).digest()[:8],
                byteorder="big"
            ) % (2**63)
            points.append(PointStruct(
                id=point_id,
                vector=embedding,
                payload={
                    "doc_id":         doc_id,
                    "url":            source,
                    "title":          title,
                    "collection":     collection,
                    "vendor":         vendor,
                    "chunk_index":    i+j,
                    "total_chunks":   total_chunks,
                    "content":        content,
                    "page":           page,
                    "has_image":      has_image,
                    "kind":           kind,
                    "access_roles":   access_roles,
                    "classification": classification,
                    "source_type":    source_type,
                    "ingested_at":    now,
                    "content_hash":   content_hash,
                }
            ))
            fts_rows.append((str(point_id), collection, source, content, title, vendor, doc_id))
        client.upsert(collection_name=collection, points=points)

        async with aiosqlite.connect(DB_PATH) as db:
            await db.executemany(
                "INSERT INTO fts(point_id, collection, url, content, title, vendor, doc_id) "
                "VALUES (?,?,?,?,?,?,?)",
                fts_rows
            )
            await db.commit()

        await asyncio.sleep(0.1)

    return text_count

# ── URL ingestion task ────────────────────────────────────────────────────────
async def ingest_url_task(
    doc_id: str,
    url: str,
    collection: str,
    vendor: str,
    access_roles: list[str],
    classification: str,
):
    now = datetime.utcnow().isoformat()

    async with aiosqlite.connect(DB_PATH) as db:
        await db.execute(
            "UPDATE documents SET status=?, updated_at=? WHERE id=?",
            ("processing", now, doc_id)
        )
        await db.commit()

    try:
        title, text = await fetch_url(url)

        if len(text) < 100:
            raise ValueError(f"Insufficient content extracted: {len(text)} chars")

        content_hash = hashlib.sha256(text.encode()).hexdigest()

        async with aiosqlite.connect(DB_PATH) as db:
            async with db.execute(
                "SELECT content_hash FROM documents WHERE id=?", (doc_id,)
            ) as cursor:
                row = await cursor.fetchone()
                if row and row[0] == content_hash:
                    await db.execute(
                        "UPDATE documents SET status=?, last_checked=?, updated_at=? WHERE id=?",
                        ("unchanged", now, now, doc_id)
                    )
                    await db.commit()
                    return

        chunk_count = await run_pipeline(
            doc_id, url, title,
            [{"page": None, "text": text, "has_image": False}], content_hash,
            collection, vendor, access_roles, classification, "url"
        )

        async with aiosqlite.connect(DB_PATH) as db:
            await db.execute("""
                UPDATE documents
                SET status=?, title=?, content_hash=?, chunk_count=?,
                    updated_at=?, last_checked=?
                WHERE id=?
            """, ("completed", title, content_hash, chunk_count, now, now, doc_id))
            await db.commit()

    except Exception as e:
        error_msg = str(e) or repr(e) or type(e).__name__
        async with aiosqlite.connect(DB_PATH) as db:
            await db.execute(
                "UPDATE documents SET status=?, error=?, updated_at=? WHERE id=?",
                ("failed", error_msg, now, doc_id)
            )
            await db.commit()
        raise

# ── Document ingestion task ───────────────────────────────────────────────────
async def ingest_document_task(
    doc_id: str,
    filename: str,
    content: bytes,
    collection: str,
    vendor: str,
    access_roles: list[str],
    classification: str,
):
    now = datetime.utcnow().isoformat()

    async with aiosqlite.connect(DB_PATH) as db:
        await db.execute(
            "UPDATE documents SET status=?, updated_at=? WHERE id=?",
            ("processing", now, doc_id)
        )
        await db.commit()

    try:
        title, segments, viewable = await extract_document(filename, content)

        full_text = "\n\n".join(s["text"] for s in segments)
        if len(full_text) < 50:
            raise ValueError(f"Insufficient content extracted: {len(full_text)} chars")

        content_hash = hashlib.sha256(full_text.encode()).hexdigest()

        async with aiosqlite.connect(DB_PATH) as db:
            async with db.execute(
                "SELECT content_hash FROM documents WHERE id=?", (doc_id,)
            ) as cursor:
                row = await cursor.fetchone()
                if row and row[0] == content_hash:
                    await db.execute(
                        "UPDATE documents SET status=?, last_checked=?, updated_at=? WHERE id=?",
                        ("unchanged", now, now, doc_id)
                    )
                    await db.commit()
                    return

        # Always keep the original upload. For an office doc we also derived a page-accurate PDF
        # (viewable) — store it too and point citations/page-images at it, since the "p.N" citations
        # index the PDF's pages and it renders inline in chat. PDFs: viewable IS content (skip).
        file_path = save_file(doc_id, filename, content, vendor)
        if viewable is not None and viewable is not content:
            stem = filename.rsplit(".", 1)[0] if "." in filename else filename
            file_path = save_file(doc_id, f"{stem}.pdf", viewable, vendor)

        chunk_count = await run_pipeline(
            doc_id, filename, title, segments, content_hash,
            collection, vendor, access_roles, classification, "document"
        )

        async with aiosqlite.connect(DB_PATH) as db:
            await db.execute("""
                UPDATE documents
                SET status=?, title=?, content_hash=?, chunk_count=?,
                    file_path=?, updated_at=?, last_checked=?
                WHERE id=?
            """, ("completed", title, content_hash, chunk_count, file_path, now, now, doc_id))
            await db.commit()

    except Exception as e:
        async with aiosqlite.connect(DB_PATH) as db:
            await db.execute(
                "UPDATE documents SET status=?, error=?, updated_at=? WHERE id=?",
                ("failed", str(e), now, doc_id)
            )
            await db.commit()
        raise

# ── Watch folder ─────────────────────────────────────────────────────────────
async def ingest_and_move(
    doc_id: str,
    filename: str,
    content: bytes,
    vendor: str,
    src_path: str,
    processed_dir: str,
):
    async with embed_semaphore:
        try:
            await ingest_document_task(doc_id, filename, content, vendor, vendor, [vendor], "public")
            if os.path.exists(src_path):
                os.rename(src_path, os.path.join(processed_dir, filename))
        except Exception as e:
            import logging
            logging.getLogger("watch_folder").error("ingest_and_move failed for %s/%s: %s", vendor, filename, e)


async def watch_folder():
    """Poll WATCH_DIR/<vendor>/ for new files and ingest them.
    Vendor subfolder name becomes both the vendor tag and Qdrant collection."""
    import logging
    logger = logging.getLogger("watch_folder")
    logger.info("Watch folder started: %s (poll every %ss)", WATCH_DIR, WATCH_POLL_INTERVAL)

    while True:
        try:
            entries = os.listdir(WATCH_DIR)
        except Exception as e:
            logger.error("Cannot list WATCH_DIR %s: %s", WATCH_DIR, e)
            await asyncio.sleep(WATCH_POLL_INTERVAL)
            continue

        for vendor in entries:
            # Skip non-directories and system/hidden folders
            vendor_dir = os.path.join(WATCH_DIR, vendor)
            if not os.path.isdir(vendor_dir) or vendor.startswith((".", "@")):
                continue

            try:
                processed_dir = os.path.join(vendor_dir, "processed")
                os.makedirs(processed_dir, exist_ok=True)

                for filename in os.listdir(vendor_dir):
                    filepath = os.path.join(vendor_dir, filename)
                    if not os.path.isfile(filepath):
                        continue
                    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
                    if ext not in SUPPORTED_EXTENSIONS:
                        continue

                    with open(filepath, "rb") as f:
                        content = f.read()
                    doc_id = hashlib.sha256(content).hexdigest()[:16]

                    async with aiosqlite.connect(DB_PATH) as db:
                        async with db.execute(
                            "SELECT status FROM documents WHERE id=?", (doc_id,)
                        ) as cursor:
                            row = await cursor.fetchone()

                    if row:
                        status = row[0]
                        if status == "completed" and os.path.exists(filepath):
                            shutil.move(filepath, os.path.join(processed_dir, filename))
                        if status in ("completed", "processing", "unchanged"):
                            continue
                        # Re-queue failed or stuck pending documents

                    logger.info("New file detected: %s/%s — queuing ingestion", vendor, filename)
                    now = datetime.utcnow().isoformat()
                    async with aiosqlite.connect(DB_PATH) as db:
                        await db.execute("""
                            INSERT OR REPLACE INTO documents
                            (id, url, collection, vendor, source_type, status, created_at, updated_at)
                            VALUES (?, ?, ?, ?, 'document', 'pending', ?, ?)
                        """, (doc_id, filename, vendor, vendor, now, now))
                        await db.commit()

                    asyncio.create_task(
                        ingest_and_move(doc_id, filename, content, vendor, filepath, processed_dir)
                    )

            except Exception as e:
                logger.error("Error processing vendor folder %s: %s", vendor, e)

        await asyncio.sleep(WATCH_POLL_INTERVAL)


# ── Deep crawl task ───────────────────────────────────────────────────────────
async def ingest_deep_task(doc_id: str, request: "DeepIngestRequest"):
    now = datetime.utcnow().isoformat()

    async with aiosqlite.connect(DB_PATH) as db:
        await db.execute(
            "UPDATE documents SET status=?, updated_at=? WHERE id=?",
            ("processing", now, doc_id)
        )
        await db.commit()

    try:
        from crawl4ai.deep_crawling import BFSDeepCrawlStrategy
        try:
            from crawl4ai.deep_crawling.filters import URLPatternFilter, FilterChain
            filter_chain = FilterChain([URLPatternFilter(patterns=[request.include_pattern])]) \
                if request.include_pattern else None
            strategy = BFSDeepCrawlStrategy(
                max_depth=request.max_depth,
                max_pages=request.max_pages,
                filter_chain=filter_chain,
            )
        except ImportError:
            strategy = BFSDeepCrawlStrategy(
                max_depth=request.max_depth,
                max_pages=request.max_pages,
            )

        base_config = dict(
            cache_mode=CacheMode.BYPASS,
            wait_until="domcontentloaded",
            page_timeout=45000,
            remove_overlay_elements=True,
            excluded_tags=["nav", "footer", "header", "aside"],
            word_count_threshold=10,
            magic=True,
        )

        async with crawl_semaphore:
            async with AsyncWebCrawler(headless=True) as crawler:
                results = await crawler.arun(
                    url=request.url,
                    config=CrawlerRunConfig(**base_config, deep_crawl_strategy=strategy),
                )

        if not isinstance(results, list):
            results = [results]

        pages_ingested = 0
        total_chunks = 0

        for result in results:
            if not result.success or not result.markdown:
                continue

            page_url = result.url
            title = (result.metadata or {}).get("title", "").strip() or page_url
            text = re.sub(r'\s+', ' ', result.markdown).strip()

            if len(text) < 100:
                continue

            page_doc_id = hashlib.sha256(page_url.encode()).hexdigest()[:16]
            content_hash = hashlib.sha256(text.encode()).hexdigest()
            page_now = datetime.utcnow().isoformat()

            async with aiosqlite.connect(DB_PATH) as db:
                await db.execute("""
                    INSERT OR REPLACE INTO documents
                    (id, url, collection, vendor, source_type, status, created_at, updated_at)
                    VALUES (?, ?, ?, ?, 'deep_crawl', 'processing', ?, ?)
                """, (page_doc_id, page_url, request.collection, request.vendor, page_now, page_now))
                await db.commit()

            try:
                chunk_count = await run_pipeline(
                    page_doc_id, page_url, title,
                    [{"page": None, "text": text, "has_image": False}], content_hash,
                    request.collection, request.vendor,
                    request.access_roles, request.classification, "deep_crawl"
                )
                pages_ingested += 1
                total_chunks += chunk_count

                async with aiosqlite.connect(DB_PATH) as db:
                    await db.execute("""
                        UPDATE documents
                        SET status=?, title=?, content_hash=?, chunk_count=?, updated_at=?, last_checked=?
                        WHERE id=?
                    """, ("completed", title, content_hash, chunk_count, page_now, page_now, page_doc_id))
                    await db.commit()

            except Exception as e:
                async with aiosqlite.connect(DB_PATH) as db:
                    await db.execute(
                        "UPDATE documents SET status=?, error=?, updated_at=? WHERE id=?",
                        ("failed", str(e), page_now, page_doc_id)
                    )
                    await db.commit()

        async with aiosqlite.connect(DB_PATH) as db:
            await db.execute("""
                UPDATE documents
                SET status=?, title=?, chunk_count=?, updated_at=?, last_checked=?
                WHERE id=?
            """, (
                "completed",
                f"Deep crawl: {request.url} ({pages_ingested} pages)",
                total_chunks, now, now, doc_id
            ))
            await db.commit()

    except Exception as e:
        async with aiosqlite.connect(DB_PATH) as db:
            await db.execute(
                "UPDATE documents SET status=?, error=?, updated_at=? WHERE id=?",
                ("failed", str(e), now, doc_id)
            )
            await db.commit()
        raise


# ── Request/Response models ───────────────────────────────────────────────────
class IngestRequest(BaseModel):
    url: str
    collection: str
    vendor: str
    access_roles: Optional[list[str]] = ["all"]
    classification: Optional[str] = "public"

class BatchIngestRequest(BaseModel):
    documents: list[IngestRequest]

class DeepIngestRequest(BaseModel):
    url: str
    collection: str
    vendor: str
    max_depth: int = 2
    max_pages: int = 30
    include_pattern: Optional[str] = None
    access_roles: Optional[list[str]] = ["all"]
    classification: Optional[str] = "public"

class CollectionCreateRequest(BaseModel):
    name: str

class MoveDocumentRequest(BaseModel):
    target_collection: str

class SetVendorRequest(BaseModel):
    vendor: str

class RenameCollectionRequest(BaseModel):
    new_name: str

# ── API Endpoints ─────────────────────────────────────────────────────────────

@app.post("/ingest/url", dependencies=_guard(INGEST_WRITE))
async def ingest_single(request: IngestRequest, background_tasks: BackgroundTasks):
    doc_id = hashlib.sha256(request.url.encode()).hexdigest()[:16]
    now = datetime.utcnow().isoformat()

    async with aiosqlite.connect(DB_PATH) as db:
        await db.execute("""
            INSERT OR REPLACE INTO documents
            (id, url, collection, vendor, source_type, status, created_at, updated_at)
            VALUES (?, ?, ?, ?, 'url', 'pending', ?, ?)
        """, (doc_id, request.url, request.collection, request.vendor, now, now))
        await db.commit()

    background_tasks.add_task(
        ingest_url_task, doc_id, request.url, request.collection,
        request.vendor, request.access_roles, request.classification
    )

    return {"doc_id": doc_id, "status": "pending", "message": f"Ingestion started for {request.url}"}

@app.post("/ingest/batch", dependencies=_guard(INGEST_WRITE))
async def ingest_batch(request: BatchIngestRequest, background_tasks: BackgroundTasks):
    results = []
    for doc in request.documents:
        doc_id = hashlib.sha256(doc.url.encode()).hexdigest()[:16]
        now = datetime.utcnow().isoformat()

        async with aiosqlite.connect(DB_PATH) as db:
            await db.execute("""
                INSERT OR REPLACE INTO documents
                (id, url, collection, vendor, source_type, status, created_at, updated_at)
                VALUES (?, ?, ?, ?, 'url', 'pending', ?, ?)
            """, (doc_id, doc.url, doc.collection, doc.vendor, now, now))
            await db.commit()

        background_tasks.add_task(
            ingest_url_task, doc_id, doc.url, doc.collection,
            doc.vendor, doc.access_roles, doc.classification
        )

        results.append({"doc_id": doc_id, "url": doc.url, "status": "pending"})

    return {"submitted": len(results), "documents": results}

@app.post("/ingest/deep", dependencies=_guard(INGEST_WRITE))
async def ingest_deep(request: DeepIngestRequest, background_tasks: BackgroundTasks):
    doc_id = hashlib.sha256(request.url.encode()).hexdigest()[:16]
    now = datetime.utcnow().isoformat()

    async with aiosqlite.connect(DB_PATH) as db:
        await db.execute("""
            INSERT OR REPLACE INTO documents
            (id, url, collection, vendor, source_type, status, created_at, updated_at)
            VALUES (?, ?, ?, ?, 'deep_crawl', 'pending', ?, ?)
        """, (doc_id, request.url, request.collection, request.vendor, now, now))
        await db.commit()

    background_tasks.add_task(ingest_deep_task, doc_id, request)

    return {"doc_id": doc_id, "status": "pending",
            "message": f"Deep crawl started for {request.url}"}


@app.post("/ingest/document", dependencies=_guard(INGEST_WRITE))
async def ingest_document(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    collection: str = Form(...),
    vendor: str = Form(...),
    access_roles: str = Form(default="all"),
    classification: str = Form(default="public"),
):
    filename = file.filename or "upload"
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext not in SUPPORTED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type: .{ext}. Supported: {', '.join(SUPPORTED_EXTENSIONS)}"
        )

    content = await file.read()
    doc_id = hashlib.sha256(content).hexdigest()[:16]
    now = datetime.utcnow().isoformat()
    roles = [r.strip() for r in access_roles.split(",")]

    async with aiosqlite.connect(DB_PATH) as db:
        await db.execute("""
            INSERT OR REPLACE INTO documents
            (id, url, collection, vendor, source_type, status, created_at, updated_at)
            VALUES (?, ?, ?, ?, 'document', 'pending', ?, ?)
        """, (doc_id, filename, collection, vendor, now, now))
        await db.commit()

    background_tasks.add_task(
        ingest_document_task, doc_id, filename, content,
        collection, vendor, roles, classification
    )

    return {"doc_id": doc_id, "status": "pending", "message": f"Ingestion started for {filename}"}

@app.get("/documents", dependencies=_guard(INGEST_READ))
async def list_documents(
    collection: Optional[str] = None,
    status: Optional[str] = None,
    source_type: Optional[str] = None,
):
    query = "SELECT * FROM documents WHERE 1=1"
    params = []

    if collection:
        query += " AND collection=?"
        params.append(collection)
    if status:
        query += " AND status=?"
        params.append(status)
    if source_type:
        query += " AND source_type=?"
        params.append(source_type)

    query += " ORDER BY updated_at DESC"

    async with aiosqlite.connect(DB_PATH) as db:
        db.row_factory = aiosqlite.Row
        async with db.execute(query, params) as cursor:
            rows = await cursor.fetchall()
            return {"documents": [dict(row) for row in rows]}

@app.get("/documents/{doc_id}", dependencies=_guard(INGEST_READ))
async def get_document(doc_id: str):
    async with aiosqlite.connect(DB_PATH) as db:
        db.row_factory = aiosqlite.Row
        async with db.execute(
            "SELECT * FROM documents WHERE id=?", (doc_id,)
        ) as cursor:
            row = await cursor.fetchone()
            if not row:
                raise HTTPException(status_code=404, detail="Document not found")
            return dict(row)

@app.get("/documents/{doc_id}/pages/{page}/image")
async def get_page_image(doc_id: str, page: int):
    """Render a single PDF page to PNG on demand from the stored source file.
    Lets the UI surface diagram/image pages so a human can read the design."""
    path = await _locate_source(doc_id)
    if not path:
        raise HTTPException(status_code=404, detail="Source file not found for this document")
    if not path.lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Page images are only available for PDF sources")

    def render() -> bytes:
        import fitz
        with fitz.open(path) as doc:
            if page < 1 or page > doc.page_count:
                raise IndexError
            return doc[page - 1].get_pixmap(dpi=OCR_DPI).tobytes("png")

    try:
        png = await asyncio.get_event_loop().run_in_executor(None, render)
    except IndexError:
        raise HTTPException(status_code=404, detail=f"Page {page} out of range")

    import io
    from fastapi.responses import StreamingResponse
    return StreamingResponse(io.BytesIO(png), media_type="image/png")

# OPEN (no auth): the chat UI links to this so a user can open the cited source. Browsers load it
# directly and can't send a bearer token — gate it at the network, like the page-image route.
@app.get("/documents/{doc_id}/file")
async def get_document_file(doc_id: str):
    """Serve a document's original uploaded file (PDF/DOCX/PPTX/txt/md) for source citations."""
    from fastapi.responses import FileResponse
    path = await _locate_source(doc_id)
    if not path or not os.path.exists(path):
        raise HTTPException(status_code=404, detail="Source file not found for this document")
    return FileResponse(path, filename=os.path.basename(path), content_disposition_type="inline")

@app.delete("/documents/{doc_id}", dependencies=_guard(DOCS_MANAGE))
async def delete_document(doc_id: str):
    async with aiosqlite.connect(DB_PATH) as db:
        db.row_factory = aiosqlite.Row
        async with db.execute(
            "SELECT * FROM documents WHERE id=?", (doc_id,)
        ) as cursor:
            row = await cursor.fetchone()
            if not row:
                raise HTTPException(status_code=404, detail="Document not found")
            doc = dict(row)

    client = get_qdrant()
    existing = [c.name for c in client.get_collections().collections]
    if doc["collection"] in existing:
        client.delete(
            collection_name=doc["collection"],
            points_selector=Filter(
                must=[FieldCondition(key="doc_id", match=MatchValue(value=doc_id))]
            )
        )

    async with aiosqlite.connect(DB_PATH) as db:
        await db.execute("DELETE FROM documents WHERE id=?", (doc_id,))
        await db.commit()

    return {"message": f"Document {doc_id} deleted successfully"}

@app.post("/documents/{doc_id}/move", dependencies=_guard(INGEST_WRITE))
async def move_document(doc_id: str, request: MoveDocumentRequest):
    """Move a document to another collection. A 'collection' is a real Qdrant
    collection, so this migrates the doc's points (carrying their vectors — no
    re-embed) and updates the SQLite + FTS5 records."""
    target = request.target_collection.strip()
    if not target:
        raise HTTPException(status_code=400, detail="target_collection is required")

    async with aiosqlite.connect(DB_PATH) as db:
        db.row_factory = aiosqlite.Row
        async with db.execute("SELECT * FROM documents WHERE id=?", (doc_id,)) as cursor:
            row = await cursor.fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="Document not found")
    source = dict(row)["collection"]
    if source == target:
        raise HTTPException(status_code=400, detail="Document is already in that collection")

    await ensure_collection(target)
    client = get_qdrant()
    existing = [c.name for c in client.get_collections().collections]
    flt = Filter(must=[FieldCondition(key="doc_id", match=MatchValue(value=doc_id))])

    moved = 0
    if source in existing:
        moved = _copy_points(client, source, target, scroll_filter=flt)
        client.delete(collection_name=source, points_selector=flt)

    now = datetime.utcnow().isoformat()
    async with aiosqlite.connect(DB_PATH) as db:
        await db.execute("UPDATE documents SET collection=?, updated_at=? WHERE id=?", (target, now, doc_id))
        await db.execute("UPDATE fts SET collection=? WHERE doc_id=?", (target, doc_id))
        await db.commit()

    return {"message": f"Moved document from '{source}' to '{target}'", "points_moved": moved}

@app.post("/documents/{doc_id}/vendor", dependencies=_guard(INGEST_WRITE))
async def set_document_vendor(doc_id: str, request: SetVendorRequest):
    """Update a document's vendor / source tag in place — across its Qdrant point payloads
    (set_payload, no re-embed and point IDs unchanged) plus the SQLite + FTS5 records."""
    new_vendor = request.vendor.strip()
    if not new_vendor:
        raise HTTPException(status_code=400, detail="vendor is required")

    async with aiosqlite.connect(DB_PATH) as db:
        db.row_factory = aiosqlite.Row
        async with db.execute("SELECT * FROM documents WHERE id=?", (doc_id,)) as cursor:
            row = await cursor.fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="Document not found")
    doc = dict(row)

    client = get_qdrant()
    existing = [c.name for c in client.get_collections().collections]
    if doc["collection"] in existing:
        client.set_payload(
            collection_name=doc["collection"],
            payload={"vendor": new_vendor},
            points=Filter(must=[FieldCondition(key="doc_id", match=MatchValue(value=doc_id))]),
        )

    now = datetime.utcnow().isoformat()
    async with aiosqlite.connect(DB_PATH) as db:
        await db.execute("UPDATE documents SET vendor=?, updated_at=? WHERE id=?", (new_vendor, now, doc_id))
        await db.execute("UPDATE fts SET vendor=? WHERE doc_id=?", (new_vendor, doc_id))
        await db.commit()

    return {"message": f"Vendor updated to '{new_vendor}'", "doc_id": doc_id, "vendor": new_vendor}

@app.get("/collections", dependencies=_guard(INGEST_READ))
async def list_collections():
    client = get_qdrant()
    return {
        "collections": [
            {"name": c.name} for c in client.get_collections().collections
        ]
    }

@app.post("/collections", dependencies=_guard(DOCS_MANAGE))
async def create_collection(request: CollectionCreateRequest):
    await ensure_collection(request.name)
    return {"message": f"Collection '{request.name}' ready"}

@app.post("/collections/{name}/rename", dependencies=_guard(DOCS_MANAGE))
async def rename_collection(name: str, request: RenameCollectionRequest):
    """Rename a collection. Qdrant has no native rename, so this creates the new
    collection, migrates every point into it (carrying vectors), drops the old
    one, and re-tags the SQLite + FTS5 records. O(points) — can be slow for large
    collections."""
    new_name = request.new_name.strip()
    if not new_name:
        raise HTTPException(status_code=400, detail="new_name is required")
    if new_name == name:
        raise HTTPException(status_code=400, detail="new_name must differ from the current name")

    client = get_qdrant()
    existing = [c.name for c in client.get_collections().collections]
    if name not in existing:
        raise HTTPException(status_code=404, detail=f"Collection '{name}' not found")
    if new_name in existing:
        raise HTTPException(status_code=409, detail=f"Collection '{new_name}' already exists")

    await ensure_collection(new_name)
    moved = _copy_points(client, name, new_name)
    client.delete_collection(collection_name=name)

    now = datetime.utcnow().isoformat()
    async with aiosqlite.connect(DB_PATH) as db:
        await db.execute("UPDATE documents SET collection=?, updated_at=? WHERE collection=?", (new_name, now, name))
        await db.execute("UPDATE fts SET collection=? WHERE collection=?", (new_name, name))
        await db.commit()

    return {"message": f"Renamed collection '{name}' to '{new_name}'", "points_moved": moved}

@app.get("/search/lexical", dependencies=_guard(INGEST_READ))
async def lexical_search(
    q: str,
    collection: Optional[str] = None,
    limit: int = 20,
):
    if not q.strip():
        return {"results": []}
    try:
        async with aiosqlite.connect(DB_PATH) as db:
            db.row_factory = aiosqlite.Row
            if collection:
                async with db.execute(
                    "SELECT point_id, collection, url, content, title, vendor "
                    "FROM fts WHERE fts MATCH ? AND collection = ? ORDER BY rank LIMIT ?",
                    (q, collection, limit)
                ) as cursor:
                    rows = await cursor.fetchall()
            else:
                async with db.execute(
                    "SELECT point_id, collection, url, content, title, vendor "
                    "FROM fts WHERE fts MATCH ? ORDER BY rank LIMIT ?",
                    (q, limit)
                ) as cursor:
                    rows = await cursor.fetchall()
        return {"results": [dict(row) for row in rows]}
    except Exception:
        return {"results": []}


@app.get("/health")
async def health():
    return {
        "status": "healthy",
        "qdrant_url": QDRANT_URL,
        "embedding_url": EMBEDDING_URL,
    }
